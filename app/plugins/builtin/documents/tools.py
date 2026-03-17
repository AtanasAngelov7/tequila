"""Documents plugin tool definitions and implementations (Sprint 13, §8.6).

All 32 document tools across PDF, PPTX, HTML-presentation, DOCX, XLSX,
CSV, data analysis and chart domains.

Blocking I/O operations are wrapped in ``asyncio.to_thread()`` so they do
not stall the event loop.

Dependencies (pip):
    PyMuPDF (fitz), pymupdf4llm, pypdf,
    python-pptx, python-docx,
    openpyxl, fpdf2, duckdb,
    matplotlib, Pillow
"""
from __future__ import annotations

import asyncio
import base64
import csv
import io
import json
import logging
import os
import textwrap
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Small helpers
# ---------------------------------------------------------------------------

def _path(p: str) -> Path:
    """Resolve *p* to an absolute Path; expand ~ and env-vars."""
    return Path(os.path.expandvars(os.path.expanduser(p))).resolve()


def _ensure_parent(p: Path) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)


# ===========================================================================
# PDF tools
# ===========================================================================

def _pdf_open_sync(path: str) -> dict[str, Any]:
    import fitz  # PyMuPDF
    doc = fitz.open(_path(path))
    meta = doc.metadata
    return {
        "path": str(_path(path)),
        "page_count": len(doc),
        "title": meta.get("title", ""),
        "author": meta.get("author", ""),
        "subject": meta.get("subject", ""),
        "creator": meta.get("creator", ""),
        "encrypted": doc.is_encrypted,
    }


async def pdf_open(path: str) -> dict[str, Any]:
    """Open a PDF and return metadata."""
    return await asyncio.to_thread(_pdf_open_sync, path)


def _pdf_read_pages_sync(path: str, start: int, end: int) -> dict[str, Any]:
    import fitz
    import pymupdf4llm  # noqa: F401 — registers text extractor
    doc = fitz.open(_path(path))
    total = len(doc)
    s = max(0, start - 1)
    e = min(total, end)
    pages: list[dict] = []
    for i in range(s, e):
        page = doc[i]
        text = page.get_text("text")
        pages.append({"page": i + 1, "text": text})
    return {"path": str(_path(path)), "total_pages": total, "pages": pages}


async def pdf_read_pages(path: str, start: int = 1, end: int = 5) -> dict[str, Any]:
    """Extract text from a page range of a PDF."""
    return await asyncio.to_thread(_pdf_read_pages_sync, path, start, end)


def _pdf_extract_tables_sync(path: str, page_numbers: list[int] | None) -> dict[str, Any]:
    import fitz
    doc = fitz.open(_path(path))
    pages_to_scan = page_numbers if page_numbers else list(range(1, len(doc) + 1))
    results: list[dict] = []
    for pn in pages_to_scan:
        if 1 <= pn <= len(doc):
            page = doc[pn - 1]
            tables = page.find_tables()
            for ti, tbl in enumerate(tables):
                rows = tbl.extract()
                results.append({"page": pn, "table_index": ti, "rows": rows})
    return {"path": str(_path(path)), "tables": results}


async def pdf_extract_tables(path: str, page_numbers: list[int] | None = None) -> dict[str, Any]:
    """Extract tables from a PDF."""
    return await asyncio.to_thread(_pdf_extract_tables_sync, path, page_numbers)


def _pdf_extract_images_sync(path: str, output_dir: str) -> dict[str, Any]:
    import fitz
    doc = fitz.open(_path(path))
    out = _path(output_dir)
    _ensure_parent(out / "x")
    saved: list[str] = []
    for i, page in enumerate(doc):
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            ext = base_image["ext"]
            img_bytes = base_image["image"]
            dest = out / f"page{i+1}_img{img_index}.{ext}"
            dest.write_bytes(img_bytes)
            saved.append(str(dest))
    return {"extracted": saved}


async def pdf_extract_images(path: str, output_dir: str) -> dict[str, Any]:
    """Extract embedded images from a PDF to *output_dir*."""
    return await asyncio.to_thread(_pdf_extract_images_sync, path, output_dir)


def _pdf_page_to_image_sync(path: str, page_number: int, dpi: int) -> dict[str, Any]:
    import fitz
    doc = fitz.open(_path(path))
    if page_number < 1 or page_number > len(doc):
        raise ValueError(f"Page {page_number} out of range (1–{len(doc)})")
    page = doc[page_number - 1]
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    clip = page.get_pixmap(matrix=mat)
    img_bytes = clip.tobytes("png")
    b64 = base64.b64encode(img_bytes).decode()
    return {"page": page_number, "format": "png", "base64": b64, "width": clip.width, "height": clip.height}


async def pdf_page_to_image(path: str, page_number: int = 1, dpi: int = 150) -> dict[str, Any]:
    """Render a PDF page to a base64-encoded PNG image."""
    return await asyncio.to_thread(_pdf_page_to_image_sync, path, page_number, dpi)


def _pdf_search_sync(path: str, query: str) -> dict[str, Any]:
    import fitz
    doc = fitz.open(_path(path))
    hits: list[dict] = []
    for i, page in enumerate(doc):
        instances = page.search_for(query)
        if instances:
            hits.append({
                "page": i + 1,
                "count": len(instances),
                "snippet": page.get_text("text", clip=instances[0])[:200] if instances else "",
            })
    return {"query": query, "total_hits": sum(h["count"] for h in hits), "pages": hits}


async def pdf_search(path: str, query: str) -> dict[str, Any]:
    """Full-text search within a PDF; return matching pages and snippets."""
    return await asyncio.to_thread(_pdf_search_sync, path, query)


def _pdf_create_sync(output_path: str, content: str, title: str) -> dict[str, Any]:
    from fpdf import FPDF  # type: ignore[import-untyped]
    pdf = FPDF()
    pdf.set_title(title)
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    for line in content.splitlines():
        pdf.multi_cell(0, 6, line)
    out = _path(output_path)
    _ensure_parent(out)
    pdf.output(str(out))
    return {"path": str(out), "pages": pdf.page}


async def pdf_create(output_path: str, content: str, title: str = "Document") -> dict[str, Any]:
    """Create a new PDF from plain text content."""
    return await asyncio.to_thread(_pdf_create_sync, output_path, content, title)


def _pdf_from_markdown_sync(output_path: str, markdown: str, title: str) -> dict[str, Any]:
    """Convert Markdown to PDF via fpdf2's support for basic markdown formatting."""
    from fpdf import FPDF  # type: ignore[import-untyped]

    pdf = FPDF()
    pdf.set_title(title)
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped.startswith("# "):
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(0, 8, stripped[2:])
            pdf.set_font("Helvetica", size=11)
        elif stripped.startswith("## "):
            pdf.set_font("Helvetica", "B", 13)
            pdf.multi_cell(0, 7, stripped[3:])
            pdf.set_font("Helvetica", size=11)
        elif stripped.startswith("### "):
            pdf.set_font("Helvetica", "B", 11)
            pdf.multi_cell(0, 6, stripped[4:])
            pdf.set_font("Helvetica", size=11)
        else:
            pdf.multi_cell(0, 6, stripped or " ")
    out = _path(output_path)
    _ensure_parent(out)
    pdf.output(str(out))
    return {"path": str(out), "pages": pdf.page}


async def pdf_from_markdown(output_path: str, markdown: str, title: str = "Document") -> dict[str, Any]:
    """Convert Markdown text to a PDF document."""
    return await asyncio.to_thread(_pdf_from_markdown_sync, output_path, markdown, title)


def _pdf_merge_sync(paths: list[str], output_path: str) -> dict[str, Any]:
    import fitz
    merged = fitz.open()
    for p in paths:
        doc = fitz.open(_path(p))
        merged.insert_pdf(doc)
    out = _path(output_path)
    _ensure_parent(out)
    merged.save(str(out))
    return {"path": str(out), "total_pages": len(merged), "sources": [str(_path(p)) for p in paths]}


async def pdf_merge(paths: list[str], output_path: str) -> dict[str, Any]:
    """Merge multiple PDF files into one."""
    return await asyncio.to_thread(_pdf_merge_sync, paths, output_path)


def _pdf_split_sync(path: str, output_dir: str, ranges: list[list[int]]) -> dict[str, Any]:
    """Split a PDF into segments defined by *ranges* [[start,end], ...]."""
    import fitz
    doc = fitz.open(_path(path))
    out_dir = _path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    produced: list[str] = []
    for idx, (s, e) in enumerate(ranges):
        new_doc = fitz.open()
        new_doc.insert_pdf(doc, from_page=s - 1, to_page=e - 1)
        dest = out_dir / f"part_{idx + 1:03d}.pdf"
        new_doc.save(str(dest))
        produced.append(str(dest))
    return {"parts": produced}


async def pdf_split(path: str, output_dir: str, ranges: list[list[int]]) -> dict[str, Any]:
    """Split a PDF into parts defined by 1-based page ranges [[start,end], ...]."""
    return await asyncio.to_thread(_pdf_split_sync, path, output_dir, ranges)


def _pdf_edit_sync(path: str, output_path: str, annotations: list[dict]) -> dict[str, Any]:
    """Annotate/edit a PDF — add text or highlight annotations per page."""
    import fitz
    doc = fitz.open(_path(path))
    for ann in annotations:
        pn = ann.get("page", 1)
        if 1 <= pn <= len(doc):
            page = doc[pn - 1]
            kind = ann.get("type", "text")
            if kind == "text":
                x, y = ann.get("x", 50), ann.get("y", 50)
                page.insert_text((x, y), ann.get("text", ""), fontsize=ann.get("fontsize", 11))
            elif kind == "highlight":
                rect = fitz.Rect(ann.get("rect", [50, 50, 200, 70]))
                page.add_highlight_annot(rect)
    out = _path(output_path)
    _ensure_parent(out)
    doc.save(str(out))
    return {"path": str(out), "annotations_applied": len(annotations)}


async def pdf_edit(path: str, output_path: str, annotations: list[dict]) -> dict[str, Any]:
    """Add text/highlight annotations to a PDF and save to a new file."""
    return await asyncio.to_thread(_pdf_edit_sync, path, output_path, annotations)


# ===========================================================================
# PPTX tools
# ===========================================================================

def _pptx_create_sync(output_path: str, title: str, slides: list[dict]) -> dict[str, Any]:
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    title_body_layout = prs.slide_layouts[1]

    first = True
    for slide_data in slides:
        if first:
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = slide_data.get("title", title)
            if hasattr(slide.placeholders[1], 'text'):
                slide.placeholders[1].text = slide_data.get("subtitle", "")
            first = False
        else:
            slide = prs.slides.add_slide(title_body_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            tf = slide.placeholders[1].text_frame
            tf.text = slide_data.get("content", "")

    out = _path(output_path)
    _ensure_parent(out)
    prs.save(str(out))
    return {"path": str(out), "slide_count": len(prs.slides)}


async def pptx_create(output_path: str, title: str, slides: list[dict]) -> dict[str, Any]:
    """Create a new PPTX presentation with given slides."""
    return await asyncio.to_thread(_pptx_create_sync, output_path, title, slides)


def _pptx_open_sync(path: str) -> dict[str, Any]:
    from pptx import Presentation  # type: ignore[import-untyped]
    prs = Presentation(str(_path(path)))
    slides: list[dict] = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slides.append({"slide": i + 1, "text": " | ".join(texts)})
    return {"path": str(_path(path)), "slide_count": len(prs.slides), "slides": slides}


async def pptx_open(path: str) -> dict[str, Any]:
    """Read a PPTX presentation and return its slide content."""
    return await asyncio.to_thread(_pptx_open_sync, path)


def _pptx_edit_sync(path: str, output_path: str, changes: list[dict]) -> dict[str, Any]:
    """Apply changes to specific slides: {slide, placeholder_index, text}."""
    from pptx import Presentation  # type: ignore[import-untyped]
    prs = Presentation(str(_path(path)))
    applied = 0
    for ch in changes:
        si = ch.get("slide", 1) - 1
        if 0 <= si < len(prs.slides):
            slide = prs.slides[si]
            pi = ch.get("placeholder_index", 0)
            if pi < len(slide.placeholders):
                slide.placeholders[pi].text = ch.get("text", "")
                applied += 1
    out = _path(output_path)
    _ensure_parent(out)
    prs.save(str(out))
    return {"path": str(out), "changes_applied": applied}


async def pptx_edit(path: str, output_path: str, changes: list[dict]) -> dict[str, Any]:
    """Edit slide placeholder text in a PPTX and save to a new file."""
    return await asyncio.to_thread(_pptx_edit_sync, path, output_path, changes)


async def pptx_list_templates() -> dict[str, Any]:
    """List available PPTX template layouts."""
    from pptx import Presentation  # type: ignore[import-untyped]
    # Default blank presentation slide layouts
    prs = Presentation()
    layouts = [{"index": i, "name": layout.name} for i, layout in enumerate(prs.slide_layouts)]
    return {"templates": layouts}


def _pptx_from_markdown_sync(output_path: str, markdown: str) -> dict[str, Any]:
    """Parse a markdown doc into slides (--- separator) and build a PPTX."""
    from pptx import Presentation  # type: ignore[import-untyped]
    prs = Presentation()
    title_layout = prs.slide_layouts[1]

    raw_slides = markdown.split("---")
    for block in raw_slides:
        lines = [l for l in block.strip().splitlines() if l.strip()]
        if not lines:
            continue
        slide = prs.slides.add_slide(title_layout)
        title_line = lines[0].lstrip("# ").strip()
        body_lines = "\n".join(lines[1:])
        slide.shapes.title.text = title_line
        tf = slide.placeholders[1].text_frame
        tf.text = body_lines

    out = _path(output_path)
    _ensure_parent(out)
    prs.save(str(out))
    return {"path": str(out), "slide_count": len(prs.slides)}


async def pptx_from_markdown(output_path: str, markdown: str) -> dict[str, Any]:
    """Convert a Markdown document (--- separated slides) into a PPTX presentation."""
    return await asyncio.to_thread(_pptx_from_markdown_sync, output_path, markdown)


# ===========================================================================
# HTML Presentation tools  (reveal.js inspired, static HTML)
# ===========================================================================

_HTML_TEMPLATE = """\
<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>{title}</title>
  <style>
    * {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{ font-family: 'Segoe UI', sans-serif; background: #1a1a2e; color: #eee; }}
    .deck {{ display: flex; flex-direction: column; }}
    section {{
      min-height: 100vh; display: flex; flex-direction: column;
      justify-content: center; align-items: flex-start;
      padding: 4rem 6rem; border-bottom: 2px solid #16213e;
    }}
    section:nth-child(odd) {{ background: #16213e; }}
    h1 {{ font-size: 2.5rem; color: #e94560; margin-bottom: 1rem; }}
    h2 {{ font-size: 1.8rem; color: #0f3460; margin-bottom: .8rem; }}
    p, li {{ font-size: 1.1rem; line-height: 1.7; color: #ccc; }}
    ul {{ padding-left: 1.5rem; }}
  </style>
</head>
<body><div class="deck">
{slides}
</div></body>
</html>
"""

_SLIDE_TEMPLATE = """\
<section>
  <h1>{title}</h1>
  <div>{body}</div>
</section>"""


def _build_html_slides(markdown: str) -> tuple[str, str]:
    """Return (title, slides_html) from markdown."""
    raw_slides = markdown.split("---")
    title = "Presentation"
    slides_html_parts: list[str] = []
    for i, block in enumerate(raw_slides):
        lines = [l for l in block.strip().splitlines() if l.strip()]
        if not lines:
            continue
        heading = lines[0].lstrip("# ").strip()
        if i == 0:
            title = heading
        body_md = "\n".join(lines[1:])
        body_html = "<br>".join(body_md.splitlines())
        slides_html_parts.append(_SLIDE_TEMPLATE.format(title=heading, body=body_html))
    return title, "\n".join(slides_html_parts)


def _html_presentation_create_sync(output_path: str, title: str, slides: list[dict]) -> dict[str, Any]:
    md_parts: list[str] = []
    for s in slides:
        md_parts.append(f"# {s.get('title', 'Slide')}\n{s.get('content', '')}")
    markdown = "\n---\n".join(md_parts)
    _, slides_html = _build_html_slides(markdown)
    html = _HTML_TEMPLATE.format(title=title, slides=slides_html)
    out = _path(output_path)
    _ensure_parent(out)
    out.write_text(html, encoding="utf-8")
    return {"path": str(out), "slide_count": len(slides)}


async def html_presentation_create(output_path: str, title: str, slides: list[dict]) -> dict[str, Any]:
    """Create a self-contained reveal.js-style HTML presentation."""
    return await asyncio.to_thread(_html_presentation_create_sync, output_path, title, slides)


def _html_from_markdown_sync(output_path: str, markdown: str) -> dict[str, Any]:
    title, slides_html = _build_html_slides(markdown)
    html = _HTML_TEMPLATE.format(title=title, slides=slides_html)
    out = _path(output_path)
    _ensure_parent(out)
    out.write_text(html, encoding="utf-8")
    slide_count = markdown.count("---") + 1
    return {"path": str(out), "slide_count": slide_count}


async def html_presentation_from_markdown(output_path: str, markdown: str) -> dict[str, Any]:
    """Convert a --- separated Markdown document into an HTML slide deck."""
    return await asyncio.to_thread(_html_from_markdown_sync, output_path, markdown)


async def html_presentation_preview(path: str) -> dict[str, Any]:
    """Return the first 3000 chars of an HTML presentation for quick preview."""
    content = await asyncio.to_thread(lambda: _path(path).read_text(encoding="utf-8"))
    return {"path": str(_path(path)), "preview": content[:3000], "total_chars": len(content)}


# ===========================================================================
# DOCX tools
# ===========================================================================

def _docx_create_sync(output_path: str, title: str, content: str) -> dict[str, Any]:
    from docx import Document  # type: ignore[import-untyped]
    doc = Document()
    doc.add_heading(title, 0)
    for para in content.splitlines():
        doc.add_paragraph(para)
    out = _path(output_path)
    _ensure_parent(out)
    doc.save(str(out))
    return {"path": str(out)}


async def docx_create(output_path: str, title: str, content: str) -> dict[str, Any]:
    """Create a new DOCX document with heading and body text."""
    return await asyncio.to_thread(_docx_create_sync, output_path, title, content)


def _docx_open_sync(path: str) -> dict[str, Any]:
    from docx import Document  # type: ignore[import-untyped]
    doc = Document(str(_path(path)))
    paragraphs = [p.text for p in doc.paragraphs]
    return {
        "path": str(_path(path)),
        "paragraph_count": len(paragraphs),
        "text": "\n".join(paragraphs),
    }


async def docx_open(path: str) -> dict[str, Any]:
    """Read a DOCX file and return its text content."""
    return await asyncio.to_thread(_docx_open_sync, path)


def _docx_edit_sync(path: str, output_path: str, replacements: list[dict]) -> dict[str, Any]:
    """Replace occurrences of find→replace in all paragraph runs."""
    from docx import Document  # type: ignore[import-untyped]
    doc = Document(str(_path(path)))
    applied = 0
    for repl in replacements:
        find = repl.get("find", "")
        replace = repl.get("replace", "")
        for para in doc.paragraphs:
            for run in para.runs:
                if find in run.text:
                    run.text = run.text.replace(find, replace)
                    applied += 1
    out = _path(output_path)
    _ensure_parent(out)
    doc.save(str(out))
    return {"path": str(out), "replacements_applied": applied}


async def docx_edit(path: str, output_path: str, replacements: list[dict]) -> dict[str, Any]:
    """Find-and-replace text in a DOCX and save to a new file."""
    return await asyncio.to_thread(_docx_edit_sync, path, output_path, replacements)


def _docx_from_markdown_sync(output_path: str, markdown: str) -> dict[str, Any]:
    from docx import Document  # type: ignore[import-untyped]
    doc = Document()
    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("- "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif stripped.startswith("1. ") or (len(stripped) > 2 and stripped[0].isdigit() and stripped[1] == "."):
            doc.add_paragraph(stripped[3:], style="List Number")
        else:
            doc.add_paragraph(stripped)
    out = _path(output_path)
    _ensure_parent(out)
    doc.save(str(out))
    return {"path": str(out)}


async def docx_from_markdown(output_path: str, markdown: str) -> dict[str, Any]:
    """Convert Markdown to a DOCX document with proper headings and lists."""
    return await asyncio.to_thread(_docx_from_markdown_sync, output_path, markdown)


# ===========================================================================
# XLSX tools
# ===========================================================================

def _xlsx_create_sync(output_path: str, sheets: list[dict]) -> dict[str, Any]:
    """Create an XLSX from a list of {name, headers, rows} dicts."""
    from openpyxl import Workbook  # type: ignore[import-untyped]
    wb = Workbook()
    wb.remove(wb.active)  # remove default blank sheet
    for sheet_data in sheets:
        ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))
        headers = sheet_data.get("headers", [])
        if headers:
            ws.append(headers)
        for row in sheet_data.get("rows", []):
            ws.append(row)
    out = _path(output_path)
    _ensure_parent(out)
    wb.save(str(out))
    return {"path": str(out), "sheet_count": len(wb.sheetnames)}


async def xlsx_create(output_path: str, sheets: list[dict]) -> dict[str, Any]:
    """Create an XLSX workbook from a list of sheet definitions."""
    return await asyncio.to_thread(_xlsx_create_sync, output_path, sheets)


def _xlsx_open_sync(path: str) -> dict[str, Any]:
    from openpyxl import load_workbook  # type: ignore[import-untyped]
    wb = load_workbook(str(_path(path)), read_only=True)
    sheets: list[dict] = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= 100:  # cap at 100 rows for safety
                break
            rows.append(list(row))
        sheets.append({"name": name, "rows": rows})
    return {"path": str(_path(path)), "sheet_count": len(wb.sheetnames), "sheets": sheets}


async def xlsx_open(path: str) -> dict[str, Any]:
    """Read an XLSX workbook and return sheet data (capped at 100 rows/sheet)."""
    return await asyncio.to_thread(_xlsx_open_sync, path)


def _xlsx_edit_sync(path: str, output_path: str, edits: list[dict]) -> dict[str, Any]:
    """Apply cell edits: {sheet, row, col, value}. 1-based row/col."""
    from openpyxl import load_workbook  # type: ignore[import-untyped]
    wb = load_workbook(str(_path(path)))
    applied = 0
    for edit in edits:
        sheet_name = edit.get("sheet", wb.sheetnames[0])
        if sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            ws.cell(row=edit.get("row", 1), column=edit.get("col", 1), value=edit.get("value"))
            applied += 1
    out = _path(output_path)
    _ensure_parent(out)
    wb.save(str(out))
    return {"path": str(out), "edits_applied": applied}


async def xlsx_edit(path: str, output_path: str, edits: list[dict]) -> dict[str, Any]:
    """Apply cell-level edits to an XLSX and save to a new file."""
    return await asyncio.to_thread(_xlsx_edit_sync, path, output_path, edits)


# ===========================================================================
# CSV tools
# ===========================================================================

def _csv_open_sync(path: str, max_rows: int) -> dict[str, Any]:
    with open(_path(path), newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows: list[list[str]] = []
        headers: list[str] = []
        for i, row in enumerate(reader):
            if i == 0:
                headers = row
            elif i <= max_rows:
                rows.append(row)
        return {
            "path": str(_path(path)),
            "headers": headers,
            "row_count_preview": len(rows),
            "rows": rows,
        }


async def csv_open(path: str, max_rows: int = 200) -> dict[str, Any]:
    """Read a CSV file, returning headers and up to *max_rows* data rows."""
    return await asyncio.to_thread(_csv_open_sync, path, max_rows)


def _csv_query_sync(path: str, sql: str) -> dict[str, Any]:
    """Run a DuckDB SQL query against a CSV file referenced via table name 'data'."""
    import duckdb  # type: ignore[import-untyped]
    abs_path = str(_path(path))
    conn = duckdb.connect()
    # Let user reference the CSV as table "data"
    adjusted_sql = sql.replace("FROM data", f"FROM read_csv_auto('{abs_path}')")
    if "FROM data" not in sql.upper():
        # Try to inject the path for common patterns
        adjusted_sql = sql.replace("data", f"read_csv_auto('{abs_path}')", 1)
    result = conn.execute(adjusted_sql).fetchall()
    description = conn.description or []
    columns = [d[0] for d in description]
    return {"columns": columns, "rows": [list(r) for r in result], "row_count": len(result)}


async def csv_query(path: str, sql: str) -> dict[str, Any]:
    """Run a DuckDB SQL query against a CSV file (reference the table as 'data')."""
    return await asyncio.to_thread(_csv_query_sync, path, sql)


def _csv_to_xlsx_sync(csv_path: str, output_path: str) -> dict[str, Any]:
    from openpyxl import Workbook  # type: ignore[import-untyped]
    wb = Workbook()
    ws = wb.active
    with open(_path(csv_path), newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        row_count = 0
        for row in reader:
            ws.append(row)
            row_count += 1
    out = _path(output_path)
    _ensure_parent(out)
    wb.save(str(out))
    return {"path": str(out), "row_count": row_count}


async def csv_to_xlsx(csv_path: str, output_path: str) -> dict[str, Any]:
    """Convert a CSV file to XLSX."""
    return await asyncio.to_thread(_csv_to_xlsx_sync, csv_path, output_path)


# ===========================================================================
# Data analysis & charting
# ===========================================================================

def _data_analyze_sync(path: str) -> dict[str, Any]:
    """Compute basic descriptive statistics on a CSV or XLSX file via DuckDB."""
    import duckdb  # type: ignore[import-untyped]
    abs_path = str(_path(path))
    conn = duckdb.connect()
    if abs_path.endswith(".csv"):
        tbl = f"read_csv_auto('{abs_path}')"
    else:
        # For xlsx, fall back to openpyxl → in-memory csv
        from openpyxl import load_workbook  # type: ignore[import-untyped]
        wb = load_workbook(abs_path, read_only=True)
        ws = wb.active
        tmp = io.StringIO()
        writer = csv.writer(tmp)
        for row in ws.iter_rows(values_only=True):
            writer.writerow(row)
        tmp.seek(0)
        tbl = f"read_csv_auto('{tmp.getvalue()}')"  # fallback path not great; use csv_path
        tmp_path = abs_path + ".tmp.csv"
        Path(tmp_path).write_text(tmp.getvalue())
        tbl = f"read_csv_auto('{tmp_path}')"

    summary_sql = f"SUMMARIZE SELECT * FROM {tbl}"
    rows = conn.execute(summary_sql).fetchall()
    desc = conn.description or []
    columns = [d[0] for d in desc]
    summary = [dict(zip(columns, row)) for row in rows]
    return {"path": abs_path, "summary": summary}


async def data_analyze(path: str) -> dict[str, Any]:
    """Run SUMMARIZE statistics on a CSV or XLSX file using DuckDB."""
    return await asyncio.to_thread(_data_analyze_sync, path)


def _data_to_chart_sync(
    path: str,
    x_column: str,
    y_columns: list[str],
    chart_type: str,
    title: str,
    output_path: str,
) -> dict[str, Any]:
    import duckdb  # type: ignore[import-untyped]
    import matplotlib  # type: ignore[import-untyped]
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt  # type: ignore[import-untyped]

    abs_path = str(_path(path))
    conn = duckdb.connect()
    if abs_path.endswith(".csv"):
        tbl = f"read_csv_auto('{abs_path}')"
    else:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
        wb = load_workbook(abs_path, read_only=True)
        ws = wb.active
        tmp_path = abs_path + ".tmp.csv"
        with open(tmp_path, "w", newline="") as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow(row)
        tbl = f"read_csv_auto('{tmp_path}')"

    all_cols = [x_column] + y_columns
    col_str = ", ".join(f'"{c}"' for c in all_cols)
    rows = conn.execute(f"SELECT {col_str} FROM {tbl} LIMIT 500").fetchall()
    x_data = [r[0] for r in rows]
    fig, ax = plt.subplots(figsize=(10, 6))
    for i, yc in enumerate(y_columns):
        y_data = [r[i + 1] for r in rows]
        if chart_type == "bar":
            ax.bar(range(len(x_data)), y_data, label=yc)
        elif chart_type == "scatter":
            ax.scatter(x_data, y_data, label=yc)
        else:
            ax.plot(x_data, y_data, label=yc)

    ax.set_title(title)
    ax.set_xlabel(x_column)
    ax.legend()
    fig.tight_layout()
    out = _path(output_path)
    _ensure_parent(out)
    fig.savefig(str(out), dpi=120)
    plt.close(fig)
    return {"path": str(out), "chart_type": chart_type}


async def data_to_chart(
    path: str,
    x_column: str,
    y_columns: list[str],
    output_path: str,
    chart_type: str = "line",
    title: str = "Chart",
) -> dict[str, Any]:
    """Generate a chart image (PNG) from CSV/XLSX data using matplotlib."""
    return await asyncio.to_thread(
        _data_to_chart_sync, path, x_column, y_columns, chart_type, title, output_path
    )


def _chart_render_sync(spec: dict[str, Any], output_path: str) -> dict[str, Any]:
    """Render a declarative chart spec to a PNG using matplotlib."""
    import matplotlib  # type: ignore[import-untyped]
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt  # type: ignore[import-untyped]

    chart_type = spec.get("type", "line")
    labels = spec.get("labels", [])
    datasets = spec.get("datasets", [])
    title = spec.get("title", "Chart")

    fig, ax = plt.subplots(figsize=(10, 6))
    for ds in datasets:
        data = ds.get("data", [])
        label = ds.get("label", "")
        if chart_type == "bar":
            ax.bar(range(len(labels)), data, label=label)
        elif chart_type == "pie":
            ax.pie(data, labels=labels)
        elif chart_type == "scatter":
            ax.scatter(range(len(data)), data, label=label)
        else:
            ax.plot(labels, data, label=label)
    ax.set_title(title)
    if chart_type != "pie":
        ax.legend()
    fig.tight_layout()
    out = _path(output_path)
    _ensure_parent(out)
    fig.savefig(str(out), dpi=120)
    plt.close(fig)
    return {"path": str(out), "chart_type": chart_type}


async def chart_render(spec: dict[str, Any], output_path: str) -> dict[str, Any]:
    """Render a declarative chart spec ({type, labels, datasets}) to a PNG image."""
    return await asyncio.to_thread(_chart_render_sync, spec, output_path)


# ===========================================================================
# Tool definition list
# ===========================================================================

DOCUMENTS_TOOLS: list[dict[str, Any]] = [
    # ── PDF ──────────────────────────────────────────────────────────────────
    {
        "name": "pdf_open",
        "fn": pdf_open,
        "description": "Open a PDF file and return its metadata (page count, title, author).",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string", "description": "Absolute or ~ path to the PDF file."}},
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "pdf_read_pages",
        "fn": pdf_read_pages,
        "description": "Extract text from a range of pages in a PDF.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "start": {"type": "integer", "default": 1, "description": "First page (1-based)."},
                "end": {"type": "integer", "default": 5, "description": "Last page (1-based, inclusive)."},
            },
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "pdf_extract_tables",
        "fn": pdf_extract_tables,
        "description": "Extract tables from specified pages of a PDF.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "page_numbers": {
                    "type": "array",
                    "items": {"type": "integer"},
                    "description": "1-based page numbers to scan. Omit to scan all pages.",
                },
            },
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "pdf_extract_images",
        "fn": pdf_extract_images,
        "description": "Extract all embedded images from a PDF and save them to a directory.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output_dir": {"type": "string", "description": "Directory to save extracted images."},
            },
            "required": ["path", "output_dir"],
        },
        "safety": "side_effect",
    },
    {
        "name": "pdf_page_to_image",
        "fn": pdf_page_to_image,
        "description": "Render a PDF page to a base64-encoded PNG image.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "page_number": {"type": "integer", "default": 1},
                "dpi": {"type": "integer", "default": 150},
            },
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "pdf_search",
        "fn": pdf_search,
        "description": "Full-text search within a PDF; returns matching pages and snippets.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "query": {"type": "string"},
            },
            "required": ["path", "query"],
        },
        "safety": "read_only",
    },
    {
        "name": "pdf_create",
        "fn": pdf_create,
        "description": "Create a new PDF file from plain text content.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "content": {"type": "string"},
                "title": {"type": "string", "default": "Document"},
            },
            "required": ["output_path", "content"],
        },
        "safety": "side_effect",
    },
    {
        "name": "pdf_from_markdown",
        "fn": pdf_from_markdown,
        "description": "Convert a Markdown string to a PDF document.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "markdown": {"type": "string"},
                "title": {"type": "string", "default": "Document"},
            },
            "required": ["output_path", "markdown"],
        },
        "safety": "side_effect",
    },
    {
        "name": "pdf_merge",
        "fn": pdf_merge,
        "description": "Merge multiple PDF files into a single output PDF.",
        "parameters": {
            "type": "object",
            "properties": {
                "paths": {"type": "array", "items": {"type": "string"}, "description": "List of source PDF paths."},
                "output_path": {"type": "string"},
            },
            "required": ["paths", "output_path"],
        },
        "safety": "side_effect",
    },
    {
        "name": "pdf_split",
        "fn": pdf_split,
        "description": "Split a PDF into parts defined by 1-based page ranges [[start,end], ...].",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output_dir": {"type": "string"},
                "ranges": {
                    "type": "array",
                    "items": {"type": "array", "items": {"type": "integer"}, "minItems": 2, "maxItems": 2},
                    "description": "List of [start, end] page ranges.",
                },
            },
            "required": ["path", "output_dir", "ranges"],
        },
        "safety": "side_effect",
    },
    {
        "name": "pdf_edit",
        "fn": pdf_edit,
        "description": "Add text or highlight annotations to a PDF and save to a new file.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output_path": {"type": "string"},
                "annotations": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "page": {"type": "integer"},
                            "type": {"type": "string", "enum": ["text", "highlight"]},
                            "text": {"type": "string"},
                            "x": {"type": "number"},
                            "y": {"type": "number"},
                        },
                    },
                },
            },
            "required": ["path", "output_path", "annotations"],
        },
        "safety": "side_effect",
    },
    # ── PPTX ─────────────────────────────────────────────────────────────────
    {
        "name": "pptx_create",
        "fn": pptx_create,
        "description": "Create a new PowerPoint (.pptx) presentation with given slides.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "title": {"type": "string"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"},
                            "content": {"type": "string"},
                        },
                    },
                },
            },
            "required": ["output_path", "title", "slides"],
        },
        "safety": "side_effect",
    },
    {
        "name": "pptx_open",
        "fn": pptx_open,
        "description": "Open a PPTX file and return its slide text content.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "pptx_edit",
        "fn": pptx_edit,
        "description": "Edit slide placeholder text in a PPTX and save to a new file.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output_path": {"type": "string"},
                "changes": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "slide": {"type": "integer"},
                            "placeholder_index": {"type": "integer"},
                            "text": {"type": "string"},
                        },
                    },
                },
            },
            "required": ["path", "output_path", "changes"],
        },
        "safety": "side_effect",
    },
    {
        "name": "pptx_list_templates",
        "fn": pptx_list_templates,
        "description": "List the available slide layout templates in a default PPTX presentation.",
        "parameters": {"type": "object", "properties": {}, "required": []},
        "safety": "read_only",
    },
    {
        "name": "pptx_from_markdown",
        "fn": pptx_from_markdown,
        "description": "Convert a --- separated Markdown document into a PPTX presentation.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "markdown": {"type": "string"},
            },
            "required": ["output_path", "markdown"],
        },
        "safety": "side_effect",
    },
    # ── HTML presentations ────────────────────────────────────────────────────
    {
        "name": "html_presentation_create",
        "fn": html_presentation_create,
        "description": "Create a self-contained HTML slide-deck (reveal.js style) from a slide list.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "title": {"type": "string"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                        },
                    },
                },
            },
            "required": ["output_path", "title", "slides"],
        },
        "safety": "side_effect",
    },
    {
        "name": "html_presentation_from_markdown",
        "fn": html_presentation_from_markdown,
        "description": "Convert a --- separated Markdown document into an HTML slide deck.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "markdown": {"type": "string"},
            },
            "required": ["output_path", "markdown"],
        },
        "safety": "side_effect",
    },
    {
        "name": "html_presentation_preview",
        "fn": html_presentation_preview,
        "description": "Return the first 3000 characters of an HTML presentation for quick preview.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
        "safety": "read_only",
    },
    # ── DOCX ─────────────────────────────────────────────────────────────────
    {
        "name": "docx_create",
        "fn": docx_create,
        "description": "Create a new DOCX document with a heading and body text.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "title": {"type": "string"},
                "content": {"type": "string"},
            },
            "required": ["output_path", "title", "content"],
        },
        "safety": "side_effect",
    },
    {
        "name": "docx_open",
        "fn": docx_open,
        "description": "Read a DOCX file and return its paragraph text.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "docx_edit",
        "fn": docx_edit,
        "description": "Find-and-replace text in a DOCX and save to a new file.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output_path": {"type": "string"},
                "replacements": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "find": {"type": "string"},
                            "replace": {"type": "string"},
                        },
                    },
                },
            },
            "required": ["path", "output_path", "replacements"],
        },
        "safety": "side_effect",
    },
    {
        "name": "docx_from_markdown",
        "fn": docx_from_markdown,
        "description": "Convert Markdown to a DOCX document with proper headings and lists.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "markdown": {"type": "string"},
            },
            "required": ["output_path", "markdown"],
        },
        "safety": "side_effect",
    },
    # ── XLSX ─────────────────────────────────────────────────────────────────
    {
        "name": "xlsx_create",
        "fn": xlsx_create,
        "description": "Create an XLSX workbook from a list of sheet definitions {name, headers, rows}.",
        "parameters": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string"},
                "sheets": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string"},
                            "headers": {"type": "array", "items": {"type": "string"}},
                            "rows": {"type": "array", "items": {"type": "array"}},
                        },
                    },
                },
            },
            "required": ["output_path", "sheets"],
        },
        "safety": "side_effect",
    },
    {
        "name": "xlsx_open",
        "fn": xlsx_open,
        "description": "Read an XLSX workbook and return sheet data (capped at 100 rows per sheet).",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "xlsx_edit",
        "fn": xlsx_edit,
        "description": "Apply cell-level edits to an XLSX and save to a new file.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "output_path": {"type": "string"},
                "edits": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "sheet": {"type": "string"},
                            "row": {"type": "integer"},
                            "col": {"type": "integer"},
                            "value": {},
                        },
                    },
                },
            },
            "required": ["path", "output_path", "edits"],
        },
        "safety": "side_effect",
    },
    # ── CSV ──────────────────────────────────────────────────────────────────
    {
        "name": "csv_open",
        "fn": csv_open,
        "description": "Read a CSV file, returning headers and up to max_rows data rows.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "max_rows": {"type": "integer", "default": 200},
            },
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "csv_query",
        "fn": csv_query,
        "description": "Run a DuckDB SQL query against a CSV file (reference the table as 'data').",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "sql": {"type": "string", "description": "SQL query; use 'data' as table name."},
            },
            "required": ["path", "sql"],
        },
        "safety": "read_only",
    },
    {
        "name": "csv_to_xlsx",
        "fn": csv_to_xlsx,
        "description": "Convert a CSV file to XLSX format.",
        "parameters": {
            "type": "object",
            "properties": {
                "csv_path": {"type": "string"},
                "output_path": {"type": "string"},
            },
            "required": ["csv_path", "output_path"],
        },
        "safety": "side_effect",
    },
    # ── Data analysis ─────────────────────────────────────────────────────────
    {
        "name": "data_analyze",
        "fn": data_analyze,
        "description": "Run SUMMARIZE descriptive statistics on a CSV or XLSX file using DuckDB.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
        "safety": "read_only",
    },
    {
        "name": "data_to_chart",
        "fn": data_to_chart,
        "description": "Generate a chart image (PNG) from CSV/XLSX data using matplotlib.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "x_column": {"type": "string"},
                "y_columns": {"type": "array", "items": {"type": "string"}},
                "output_path": {"type": "string"},
                "chart_type": {"type": "string", "enum": ["line", "bar", "scatter"], "default": "line"},
                "title": {"type": "string", "default": "Chart"},
            },
            "required": ["path", "x_column", "y_columns", "output_path"],
        },
        "safety": "side_effect",
    },
    {
        "name": "chart_render",
        "fn": chart_render,
        "description": "Render a declarative chart spec {type, title, labels, datasets} to a PNG image.",
        "parameters": {
            "type": "object",
            "properties": {
                "spec": {
                    "type": "object",
                    "properties": {
                        "type": {"type": "string", "enum": ["line", "bar", "pie", "scatter"], "default": "line"},
                        "title": {"type": "string"},
                        "labels": {"type": "array", "items": {"type": "string"}},
                        "datasets": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "label": {"type": "string"},
                                    "data": {"type": "array", "items": {"type": "number"}},
                                },
                            },
                        },
                    },
                },
                "output_path": {"type": "string"},
            },
            "required": ["spec", "output_path"],
        },
        "safety": "side_effect",
    },
]

# ── Convenience: name → fn mapping ───────────────────────────────────────────

TOOL_FN_MAP: dict[str, Any] = {t["name"]: t["fn"] for t in DOCUMENTS_TOOLS}
